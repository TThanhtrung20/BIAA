"""Tạo file Office: Word (.docx), Excel (.xlsx), PowerPoint (.pptx).

Các hàm ở đây nhận dữ liệu có cấu trúc (dict/list) - chính là dạng mà bộ não
LLM sẽ sinh ra - và xuất ra file thật. Mọi hàm trả về đường dẫn file đã tạo.

File mặc định lưu trong ~/Documents nếu người dùng không nói rõ nơi lưu.
"""
from __future__ import annotations

import os
from datetime import datetime

DEFAULT_DIR = os.path.expanduser("~/Documents")


def slug_filename(title: str, ext: str) -> str:
    """Tạo tên file an toàn từ tiêu đề, kèm mốc thời gian tránh trùng.

    Giữ chữ tiếng Việt có dấu, chỉ bỏ ký tự không hợp lệ cho tên file.
    """
    import re

    base = (title or "tai_lieu").strip()
    base = re.sub(r'[\\/:*?"<>|]+', "", base)   # bỏ ký tự cấm trong tên file
    base = re.sub(r"\s+", "_", base)
    base = base[:50].strip("_") or "tai_lieu"
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{base}_{stamp}{ext}"


def _resolve_path(path: str | None, default_name: str, ext: str) -> str:
    """Chuẩn hoá đường dẫn lưu file.

    - Không có path -> tự đặt tên kèm thời gian, lưu trong ~/Documents.
    - Chỉ có tên file -> lưu trong ~/Documents.
    - Thiếu đuôi mở rộng -> tự thêm.
    - Đảm bảo thư mục cha tồn tại.
    """
    if not path:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        os.makedirs(DEFAULT_DIR, exist_ok=True)
        return os.path.join(DEFAULT_DIR, f"{default_name}_{stamp}{ext}")

    path = os.path.expanduser(str(path))
    if not path.lower().endswith(ext):
        path += ext
    if os.path.dirname(path) == "":
        path = os.path.join(DEFAULT_DIR, path)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    return path


# --------------------------------------------------------------------------- #
# WORD (.docx)
# --------------------------------------------------------------------------- #
def create_word(title: str | None = None, blocks: list | None = None,
                path: str | None = None) -> str:
    """Tạo tài liệu Word.

    blocks: danh sách khối nội dung. Mỗi phần tử là chuỗi (đoạn văn thường)
    hoặc dict:
      {"type": "heading", "text": "...", "level": 1}
      {"type": "paragraph", "text": "..."}
      {"type": "bullet", "text": "..."}
      {"type": "bullets", "items": ["...", "..."]}
      {"type": "number", "text": "..."}
    """
    from docx import Document

    doc = Document()
    if title:
        doc.add_heading(str(title), level=0)

    for block in (blocks or []):
        if isinstance(block, str):
            doc.add_paragraph(block)
            continue
        if not isinstance(block, dict):
            continue
        btype = str(block.get("type", "paragraph")).lower()
        text = str(block.get("text", ""))
        if btype == "heading":
            level = int(block.get("level", 1) or 1)
            doc.add_heading(text, level=max(1, min(level, 9)))
        elif btype == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        elif btype == "number":
            doc.add_paragraph(text, style="List Number")
        elif btype == "bullets":
            for item in block.get("items", []):
                doc.add_paragraph(str(item), style="List Bullet")
        else:
            doc.add_paragraph(text)

    out = _resolve_path(path, "tai_lieu", ".docx")
    doc.save(out)
    return out


# --------------------------------------------------------------------------- #
# EXCEL (.xlsx)
# --------------------------------------------------------------------------- #
def create_excel(headers: list | None = None, rows: list | None = None,
                 path: str | None = None, sheet_name: str = "Sheet1",
                 sheets: list | None = None) -> str:
    """Tạo bảng tính Excel.

    Cách 1 (một sheet): truyền headers + rows.
    Cách 2 (nhiều sheet): truyền sheets=[{"name","headers","rows"}, ...].
    """
    from openpyxl import Workbook

    wb = Workbook()
    if sheets:
        first = True
        for sh in sheets:
            ws = wb.active if first else wb.create_sheet()
            ws.title = str(sh.get("name", "Sheet"))[:31]
            _fill_sheet(ws, sh.get("headers"), sh.get("rows", []))
            first = False
    else:
        ws = wb.active
        ws.title = str(sheet_name)[:31]
        _fill_sheet(ws, headers, rows or [])

    out = _resolve_path(path, "bang_tinh", ".xlsx")
    wb.save(out)
    return out


def _fill_sheet(ws, headers, rows) -> None:
    from openpyxl.styles import Font

    if headers:
        ws.append(list(headers))
        for cell in ws[1]:
            cell.font = Font(bold=True)
    for row in (rows or []):
        ws.append(list(row) if isinstance(row, (list, tuple)) else [row])

    # Tự canh độ rộng cột theo nội dung dài nhất (tối đa 60)
    for column_cells in ws.columns:
        length = max((len(str(c.value)) for c in column_cells if c.value is not None),
                     default=0)
        ws.column_dimensions[column_cells[0].column_letter].width = min(length + 2, 60)


# --------------------------------------------------------------------------- #
# POWERPOINT (.pptx)
# --------------------------------------------------------------------------- #
def create_powerpoint(title: str | None = None, slides: list | None = None,
                      path: str | None = None, subtitle: str | None = None) -> str:
    """Tạo bản trình chiếu PowerPoint.

    slides: danh sách slide nội dung, mỗi slide là dict:
      {"title": "...", "bullets": ["...", "..."]}
    Nếu có `title` thì slide đầu tiên là slide tiêu đề.
    """
    from pptx import Presentation

    prs = Presentation()

    if title:
        layout = prs.slide_layouts[0]  # Title slide
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = str(title)
        if subtitle and len(s.placeholders) > 1:
            s.placeholders[1].text = str(subtitle)

    for slide in (slides or []):
        layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = str(slide.get("title", ""))
        bullets = slide.get("bullets", [])
        body = s.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                body.text = str(bullet)
            else:
                para = body.add_paragraph()
                para.text = str(bullet)

    out = _resolve_path(path, "trinh_chieu", ".pptx")
    prs.save(out)
    return out
